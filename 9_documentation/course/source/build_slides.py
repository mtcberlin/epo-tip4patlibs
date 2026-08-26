#!/usr/bin/env python3
"""Render the workshop deck from slides.yaml.

Three slides per module — Introduction, Working through, Outcome — in two
45-minute blocks, plus an opening and a closing slide: 20 in all. The module
documents next to this script stay the source of truth; slides.yaml is their
condensed form, and this file only lays it out.

Screenshots: drop a PNG into `shots/` named after the module (`01.png`, `03.png`,
…) and the "Working through" slide uses it. Until then the slide carries a framed
placeholder naming the shot it wants — and whether taking it needs a TIP run,
because modules 1 and 5 ship with cleared outputs on purpose.

    uv run --with python-pptx --with pyyaml --with pillow python build_slides.py
"""

from __future__ import annotations

import sys
from pathlib import Path

import yaml
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
# NOTE: this file has been opened and saved in PowerPoint since it was last
# generated, so it may carry hand edits that exist nowhere else. Re-running this
# script overwrites them. Carry a content change into slides.yaml first.
OUT = HERE.parent / "TIP4PATLIBS_1_Workshop_v1.pptx"
SHOTS = HERE / "shots"

# The course red, as worn by every notebook header — the deck sits next to those
# on screen. The handouts use the mtc accent instead; one constant either way.
ACCENT = RGBColor(0xBE, 0x0F, 0x05)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
FAINT = RGBColor(0x94, 0xA3, 0xB8)
PANEL = RGBColor(0xF8, 0xFA, 0xFC)
TINT = RGBColor(0xFD, 0xF2, 0xF2)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)                       # page margin
FONT = "Arial"


# --------------------------------------------------------------------------- #
# small typesetting helpers
# --------------------------------------------------------------------------- #
def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size=14, bold=False, color=INK, space_after=6,
         italic=False, align=PP_ALIGN.LEFT, first=False, line=1.25):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    f = r.font
    f.name, f.size, f.bold, f.italic = FONT, Pt(size), bold, italic
    f.color.rgb = color
    return p


def box(slide, x, y, w, h, *, fill=PANEL, outline=None, radius=0.04,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    if hasattr(sh, "adjustments") and len(sh.adjustments):
        sh.adjustments[0] = radius
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if outline is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = outline
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    return sh


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def eyebrow(slide, left_text, right_text=None):
    tf = textbox(slide, M, Inches(0.5), W - 2 * M, Inches(0.3))
    para(tf, left_text.upper(), size=11, bold=True, color=ACCENT, first=True,
         space_after=0)
    if right_text:
        tf2 = textbox(slide, W - M - Inches(4.5), Inches(0.5), Inches(4.5), Inches(0.3))
        para(tf2, right_text, size=11, color=FAINT, align=PP_ALIGN.RIGHT,
             first=True, space_after=0)


def headline(slide, text, y=Inches(0.95), size=30):
    tf = textbox(slide, M, y, W - 2 * M, Inches(1.0))
    para(tf, text, size=size, bold=True, color=INK, first=True, space_after=0,
         line=1.1)


def footer(slide, n, total, note=""):
    tf = textbox(slide, M, H - Inches(0.62), W - 2 * M, Inches(0.3))
    para(tf, note, size=9, color=FAINT, first=True, space_after=0)
    tf2 = textbox(slide, W - M - Inches(3), H - Inches(0.62), Inches(3), Inches(0.3))
    para(tf2, f"TIP4PATLIBS · {n} / {total}", size=9, color=FAINT,
         align=PP_ALIGN.RIGHT, first=True, space_after=0)


# --------------------------------------------------------------------------- #
# the chain diagram, drawn rather than pictured
# --------------------------------------------------------------------------- #
def chain(slide, cfg, top):
    colw, gap = Inches(4.6), Inches(1.9)
    x_skills = M
    x_apps = M + colw + gap
    cellh, cellgap = Inches(0.76), Inches(0.15)
    left, right = cfg["left"], cfg["right"]
    # the two groups are drawn the same height so their frames align, even when one
    # side holds a single claim and the other a three-rung ladder
    tallest = max(len(left), len(right))
    grouph = Inches(0.40) + tallest * cellh + (tallest - 1) * cellgap + Inches(0.22)

    for x, label, rows in ((x_skills, cfg["left_label"], left),
                           (x_apps, cfg["right_label"], right)):
        # centre a short column inside the shared frame
        pad = (tallest - len(rows)) * (cellh + cellgap) / 2
        box(slide, x, top, colw, grouph, fill=None, outline=LINE, radius=0.03)
        tf = textbox(slide, x, top + Inches(0.14), colw, Inches(0.26))
        para(tf, label, size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
             first=True, space_after=0)
        for i, (num, title, sub) in enumerate(rows):
            cy = top + Inches(0.5) + pad + i * (cellh + cellgap)
            box(slide, x + Inches(0.22), cy, colw - Inches(0.44), cellh,
                fill=TINT, outline=ACCENT, radius=0.12)
            tf = textbox(slide, x + Inches(0.42), cy + Inches(0.11),
                         colw - Inches(0.84), cellh - Inches(0.16))
            para(tf, f"{num}   {title}" if num else title, size=13, bold=True,
                 first=True, space_after=1)
            para(tf, sub, size=10, color=MUTED, space_after=0)

    ax = x_skills + colw
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax + Inches(0.35),
                                   top + grouph / 2 - Inches(0.16),
                                   gap - Inches(0.7), Inches(0.32))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    tf = textbox(slide, ax + Inches(0.1), top + grouph / 2 - Inches(0.95),
                 gap - Inches(0.2), Inches(0.7))
    para(tf, cfg["arrow"], size=9, color=MUTED,
         align=PP_ALIGN.CENTER, first=True, space_after=0)


# --------------------------------------------------------------------------- #
# the five slide kinds
# --------------------------------------------------------------------------- #
def slide_title(prs, d, total):
    s = blank(prs)
    box(s, Emu(0), Emu(0), W, Inches(0.16), fill=ACCENT, radius=0,
        shape=MSO_SHAPE.RECTANGLE)
    dk = d["deck"]
    tf = textbox(s, M, Inches(0.62), W - 2 * M, Inches(1.5))
    para(tf, dk["title"], size=42, bold=True, color=ACCENT, first=True,
         space_after=4, line=1.0)
    para(tf, dk["subtitle"], size=16, color=INK, space_after=2)
    para(tf, f'{dk["occasion"]}   ·   {dk["author"]}', size=12, color=MUTED)

    goal = box(s, M, Inches(2.25), W - 2 * M, Inches(1.15), fill=TINT,
               outline=None, radius=0.06)
    tf = goal.text_frame
    tf.margin_left = tf.margin_right = Inches(0.28)
    tf.margin_top = tf.margin_bottom = Inches(0.16)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, dk["goal"], size=13, italic=True, color=INK, first=True,
         space_after=0, line=1.3)

    chain(s, dk["chain"], Inches(3.62))
    footer(s, 1, total, dk["chain"]["note"])


def slide_intro(prs, mod, blk, n, total):
    s = blank(prs)
    eyebrow(s, f'{blk} · Module {mod["n"]}', mod.get("credit", ""))
    headline(s, f'Module {mod["n"]} — {mod["title"]}')

    q = box(s, M, Inches(1.95), W - 2 * M, Inches(1.5), fill=TINT, radius=0.06)
    tf = q.text_frame
    tf.margin_left = tf.margin_right = Inches(0.34)
    tf.margin_top = tf.margin_bottom = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, f'"{mod["question"].strip()}"', size=18, italic=True, color=ACCENT,
         first=True, space_after=0, line=1.25)

    tf = textbox(s, M, Inches(3.75), W - 2 * M, Inches(0.3))
    para(tf, "Why the obvious answer fails", size=13, bold=True, color=ACCENT,
         first=True, space_after=0)

    y = Inches(4.2)
    for line in mod["fails"]:
        box(s, M, y + Inches(0.16), Inches(0.1), Inches(0.1), fill=ACCENT,
            radius=0.5, shape=MSO_SHAPE.OVAL)
        tf = textbox(s, M + Inches(0.36), y, W - 2 * M - Inches(0.36), Inches(0.55))
        para(tf, line, size=14, color=INK, first=True, space_after=0, line=1.25)
        y += Inches(0.62)

    footer(s, n, total, mod["folder"])


def slide_work(prs, mod, blk, n, total):
    s = blank(prs)
    eyebrow(s, f'{blk} · Module {mod["n"]}',
            mod.get("minutes", "Working through"))
    headline(s, mod["title"], size=26)

    shot_w, shot_h = Inches(7.0), Inches(4.35)
    x, y = M, Inches(1.9)
    png = SHOTS / f'{mod["n"]:02d}.png'
    if png.exists():
        # fit inside the frame and centre — a portrait shot scaled to the frame's
        # width would run off the bottom of the slide
        with Image.open(png) as _im:
            iw, ih = _im.size
        scale = min(shot_w / iw, shot_h / ih)
        pw, ph = int(iw * scale), int(ih * scale)
        s.shapes.add_picture(str(png), x + (shot_w - pw) // 2,
                             y + (shot_h - ph) // 2, width=pw, height=ph)
    else:
        avail = mod["shot"]["available"]
        tag = {True: "available offline — notebook ships executed",
               "partial": "partly available — some cells ship cleared",
               False: "NEEDS A TIP RUN — notebook ships with cleared outputs"}[avail]
        ph = box(s, x, y, shot_w, shot_h, fill=PANEL, outline=LINE, radius=0.02)
        tf = ph.text_frame
        tf.margin_left = tf.margin_right = Inches(0.4)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, "SCREENSHOT", size=11, bold=True, color=FAINT,
             align=PP_ALIGN.CENTER, first=True, space_after=10)
        para(tf, mod["shot"]["what"], size=15, color=INK, align=PP_ALIGN.CENTER,
             space_after=10, line=1.25)
        para(tf, mod["shot"]["file"], size=10, color=MUTED, align=PP_ALIGN.CENTER,
             space_after=8)
        para(tf, tag, size=10, bold=(avail is False),
             color=ACCENT if avail is False else FAINT, align=PP_ALIGN.CENTER,
             space_after=0)

    rx = M + shot_w + Inches(0.5)
    rw = W - M - rx
    yy = y
    for i, (what, detail, mins) in enumerate(mod["steps"], 1):
        box(s, rx, yy, Inches(0.34), Inches(0.34), fill=ACCENT, radius=0.2)
        tf = textbox(s, rx, yy + Inches(0.05), Inches(0.34), Inches(0.3))
        para(tf, str(i), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             first=True, space_after=0)
        tf = textbox(s, rx + Inches(0.5), yy - Inches(0.02), rw - Inches(0.5),
                     Inches(1.2))
        para(tf, what, size=15, bold=True, first=True, space_after=3)
        para(tf, detail, size=11, color=MUTED, space_after=3, line=1.2)
        para(tf, mins, size=10, bold=True, color=ACCENT, space_after=0)
        yy += Inches(1.45)

    footer(s, n, total, mod["folder"])


def slide_outcome(prs, mod, blk, n, total):
    s = blank(prs)
    eyebrow(s, f'{blk} · Module {mod["n"]}', "Outcome")
    headline(s, "What now exists", size=28)

    y = Inches(1.95)
    for line in mod["artifact"]:
        b = box(s, M, y, W - 2 * M, Inches(0.78), fill=PANEL, radius=0.05)
        tf = b.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, line, size=14, color=INK, first=True, space_after=0, line=1.2)
        y += Inches(0.92)

    tf = textbox(s, M, y + Inches(0.15), W - 2 * M, Inches(0.3))
    para(tf, "The one sentence to remember", size=12, bold=True, color=ACCENT,
         first=True, space_after=0)
    t = box(s, M, y + Inches(0.6), W - 2 * M, Inches(1.15), fill=TINT, radius=0.05)
    t.line.color.rgb = ACCENT
    t.line.width = Pt(1)
    tf = t.text_frame
    tf.margin_left = tf.margin_right = Inches(0.34)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, mod["takeaway"].strip(), size=16, bold=True, color=ACCENT,
         first=True, space_after=0, line=1.25)

    footer(s, n, total, mod["folder"])


def slide_closing(prs, d, total):
    s = blank(prs)
    c = d["closing"]
    box(s, Emu(0), Emu(0), W, Inches(0.16), fill=ACCENT, radius=0,
        shape=MSO_SHAPE.RECTANGLE)
    headline(s, c["headline"], y=Inches(0.9), size=34)

    y = Inches(2.2)
    for line in c["points"]:
        box(s, M, y + Inches(0.2), Inches(0.12), Inches(0.12), fill=ACCENT,
            radius=0.5, shape=MSO_SHAPE.OVAL)
        tf = textbox(s, M + Inches(0.42), y, W - 2 * M - Inches(0.42), Inches(0.7))
        para(tf, line, size=17, color=INK, first=True, space_after=0, line=1.25)
        y += Inches(0.95)

    f = box(s, M, Inches(5.4), W - 2 * M, Inches(1.1), fill=TINT, radius=0.05)
    tf = f.text_frame
    tf.margin_left = tf.margin_right = Inches(0.34)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, c["handouts"], size=14, bold=True, color=ACCENT, first=True,
         space_after=3)
    para(tf, c["repo"], size=12, color=MUTED, space_after=0)

    footer(s, total, total, d["deck"]["author"])


# --------------------------------------------------------------------------- #
def main() -> int:
    d = yaml.safe_load((HERE / "slides.yaml").read_text(encoding="utf-8"))
    mods = {m["n"]: m for m in d["modules"]}
    total = 2 + 3 * len(mods)

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    slide_title(prs, d, total)
    n = 2
    for blk in d["blocks"]:
        for num in blk["modules"]:
            mod = mods[num]
            slide_intro(prs, mod, blk["name"], n, total); n += 1
            slide_work(prs, mod, blk["name"], n, total); n += 1
            slide_outcome(prs, mod, blk["name"], n, total); n += 1
    slide_closing(prs, d, total)

    prs.save(OUT)
    missing = [m["n"] for m in d["modules"] if not (SHOTS / f'{m["n"]:02d}.png').exists()]
    print(f"   {OUT.relative_to(HERE.parent.parent)}  "
          f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
          f"{OUT.stat().st_size // 1024} kB)")
    if missing:
        needs_tip = [m["n"] for m in d["modules"]
                     if m["shot"]["available"] is False]
        print(f"   screenshots still placeholders: {missing}")
        print(f"   of those, needing a TIP run:    {needs_tip}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
